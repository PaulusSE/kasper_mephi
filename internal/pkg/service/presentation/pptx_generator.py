import sys
import json
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Чтение данных JSON из stdin
input_data = json.load(sys.stdin)

# Извлечение данных из JSON
current_semester = input_data['current_semester']
full_name = input_data['full_name']
supervisor_name = input_data['supervisor_name']
education_direction = input_data['education_direction']
education_profile = input_data['education_profile']
enrollment_date = input_data['enrollment_date']
specialty = input_data['specialty']
training_year_fgos = input_data['training_year_fgos']
candidate_exams = input_data['candidate_exams']
category = input_data['category']
topic = input_data['topic']
report_period_work = input_data['report_period_work']
scientific_obj = input_data['scientific_obj']
scientific_subj = input_data['scientific_subj']
mentor_rate = input_data['mentor_rate']
progress_percents = input_data['progress_percents']
progress_descriptions = input_data['progress_descriptions']
publications = input_data['publications']
all_publications = input_data['all_publications']
pedagogical_data = input_data['pedagogical_data']
report_other_achievments = input_data['report_other_achievments']
pedagogical_data_all = input_data['pedagogical_data_all']
next_semester_plan = input_data['next_semester_plan']


class SlideGenerator:
    def __init__(self):
        self.prs = Presentation()

    def add_slide(self, layout_index=5):
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        return slide

    def add_text_to_slide(self, slide, text, left, top, width, height, font_size=12, bold=False, alignment=PP_ALIGN.LEFT, color=(0, 0, 0)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        if not tf.text:
            tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        p.font.color.rgb = RGBColor(color[0], color[1], color[2])

    def add_image_to_slide(self, slide, image_path, left, top, width, height):
        slide.shapes.add_picture(image_path, left, top, width, height)

    def add_table_to_slide(self, slide, data, left, top, width, height, font_size=8):
        rows = len(data)
        cols = max(len(row) for row in data)  # Находим максимальное количество столбцов в строках
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for i, row in enumerate(data):
            for j, (text, rowspan, colspan) in enumerate(row):
                if text is not None:
                    if isinstance(text, list):  # Проверяем, является ли текст списком
                        text = ", ".join(text)  # Преобразуем список в строку, разделенную запятыми
                    cell = table.cell(i, j)
                    cell.text = text
                    if rowspan > 1 or colspan > 1:
                        cell.merge(table.cell(i + rowspan - 1, j + colspan - 1))

        # Устанавливаем размер шрифта для всей таблицы
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                        paragraph.alignment = PP_ALIGN.LEFT  # Выравнивание текста по левому краю
        return table

    def add_decorative_elements(self, slide):
        shapes = slide.shapes
        line_top = shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3)
        )
        line_top.fill.solid()
        line_top.fill.fore_color.rgb = RGBColor(0, 112, 192)

        line_bottom = shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(10), Inches(0.3)
        )
        line_bottom.fill.solid()
        line_bottom.fill.fore_color.rgb = RGBColor(0, 112, 192)

    def add_background(self, slide, color=(245, 245, 245)):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])

    def create_presentation(self, slides_data):
        for i, slide_data in enumerate(slides_data):
            slide = self.add_slide()
            self.add_background(slide)
            self.add_decorative_elements(slide)
            for element in slide_data:
                if element['type'] == 'text':
                    self.add_text_to_slide(slide, element['content'], element['left'], element['top'], element['width'], element['height'],
                                          element.get('font_size', 12), element.get('bold', False), element.get('alignment', PP_ALIGN.LEFT),
                                          element.get('color', (0, 0, 0)))
                elif element['type'] == 'image':
                    self.add_image_to_slide(slide, element['image_path'], element['left'], element['top'], element['width'], element['height'])
                elif element['type'] == 'table':
                    table = self.add_table_to_slide(slide, element['data'], element['left'], element['top'], element['width'], element['height'], element.get('font_size', 8))
                    if 'column_widths' in element:
                        set_column_widths(table, element['column_widths'])

    def save_presentation(self, filename):
        self.prs.save(filename)

# def build_table_disser_work_data(topic, progress_percents, progress_descriptions):
#     table_data = [
#         [("Тема:", 1, 1), (topic, 1, 3)],
#         [("Прогресс (по рукописи и разработке ПО)", 8, 1)],
#         # [(),("№ семестра", 1, 1), ("%", 1, 1), ("Что конкретно сделано за семестр", 1, 3)],
#     ]

#     semesters = [
#         "1-ый семестр", "2-ой семестр", "3-ий семестр", "4-ый семестр",
#         "5-ый семестр", "6-ой семестр", "7-ой семестр", "8-ой семестр"
#     ]

#     for i in range(len(semesters)):
#         percent = progress_percents[i] if i < len(progress_percents) else ""
#         description = progress_descriptions[i] if i < len(progress_descriptions) else ""
#         table_data.append([("", 1, 1), (semesters[i], 1, 1), (percent, 1, 1), (description, 1, 1)])

#     return table_data

def build_table_disser_work_data(topic, progress_percents, progress_descriptions):
    table_data = [
        [("Тема:", 2, 1), (topic, 2, 3)],  # Объединение строк для "Тема:"
        [("", 1, 1), ("", 1, 1), ("", 1, 1), ("", 1, 1)],  # Пустые строки для объединённых ячеек
        [("Прогресс (по рукописи и разработке ПО)", 9, 1), ("По семестрам", 1, 3)],
        # [("", 1, 1),("По семестрам", 1, 3)],
    ]

    semesters = [
        "1-ый семестр", "2-ой семестр", "3-ий семестр", "4-ый семестр",
        "5-ый семестр", "6-ой семестр", "7-ой семестр", "8-ой семестр"
    ]

    for i in range(len(semesters)):
        percent = progress_percents[i] if i < len(progress_percents) else ""
        description = progress_descriptions[i] if i < len(progress_descriptions) else ""
        table_data.append([("", 1, 0), (semesters[i], 1, 0), (percent, 1, 0), (description, 1, 0)])  # Добавляем данные без объединения ячеек

    return table_data

def build_table_publications_data(publications):
    table_data = [
        [
            ("№ п/п", 1, 1),
            ("Наименование работы", 1, 1),
            ("СТАТУС (ВАК, РИНЦ, SCOPUS, WoS)", 1, 1),
            ("Импактфактор издания", 1, 1),
            ("Выходные данные", 1, 1),
            ("Объем в стр", 1, 1),
            ("Соавторы", 1, 1)
        ]
    ]

    for i, publication in enumerate(publications):
        row = [
            (str(i + 1), 1, 1),  # № п/п
            (publication[0], 1, 1),  # Наименование работы
            (publication[1], 1, 1),  # СТАТУС
            (publication[2], 1, 1),  # Импактфактор издания
            (publication[3], 1, 1),  # Выходные данные
            (publication[4], 1, 1),  # Объем в стр
            (publication[5], 1, 1)  # Соавторы
        ]
        table_data.append(row)

    column_widths = [
        Inches(0.5),  # № п/п
        Inches(2.25),  # Наименование работы
        Inches(0.8),  # СТАТУС
        Inches(0.8),  # Импактфактор издания
        Inches(2.5),  # Выходные данные
        Inches(0.7),  # Объем в стр
        Inches(1.2)  # Соавторы
    ]

    return (table_data, column_widths)

def set_column_widths(table, column_widths):
    for i, width in enumerate(column_widths):
        table.columns[i].width = width

def build_table_publ_all(data):
    headers = [
        "№ семестра", "ВАК", "WoS / Scopus", "РИНЦ",
        "Тезисы в сборниках конфер.", "Очное участие в конфер.",
        "Свидетел. о регистраци и ЭВМ", "ВСЕГО за семестр"
    ]

    rows = [
        "До поступления", "1-ый семестр", "2-ой семестр", "3-ий семестр", "4-ый семестр",
        "5-ый семестр", "6-ой семестр", "7-ой семестр", "8-ой семестр", "ТЕКУЩИЙ ИТОГ"
    ]

    # Инициализация таблицы с заголовками
    table_data = [[(header, 1, 1) for header in headers]]  # Заголовки столбцов

    # Заполнение данных по семестрам и подсчет суммы за семестр
    for i, row in enumerate(data):
        row_data = [(rows[i], 1, 1)]  # Первый столбец - название строки
        total_for_semester = 0

        for j in range(6):  # Ожидаем 6 столбцов данных
            value = row[j] if j < len(row) and row[j] != '' else ''
            if value != '':
                total_for_semester += int(value)  # Суммируем по строкам
            row_data.append((str(value), 1, 1))

        row_data.append((str(total_for_semester) if total_for_semester != 0 else '', 1, 1))  # Добавление суммы за семестр
        table_data.append(row_data)

    # Подсчет итога по столбцам
    total_row = [("ТЕКУЩИЙ ИТОГ", 1, 1)]
    for col in range(1, len(headers)):
        total_for_column = 0
        for row in table_data[1:]:
            value = int(row[col][0]) if row[col][0] != '' else 0
            total_for_column += value

        total_row.append((str(total_for_column) if total_for_column != 0 else '', 1, 1))

    table_data.append(total_row)

    return table_data

def build_table_pedagogical_work(data):
    headers = [
        "Семестр", "Группы", "Основной преподаватель",
        "Тип занятий", "Кол-во часов за весь семестр"
    ]

    # Инициализация таблицы с заголовками
    table_data = [[(header, 1, 1) for header in headers]]  # Заголовки столбцов

    # Заполнение данных
    for row in data:
        row_data = []
        for cell in row:
            value = cell if cell != '' else ''  # Если значение пустое, то заменить на пустую строку
            row_data.append((str(value), 1, 1))
        table_data.append(row_data)

    return table_data

def build_table_pedagogical_work_all(data):
    headers = [
        "№ семестра", "Кол-во аудит. часов", "Иная педагогическая практика"
    ]

    rows = [
        "1-ый семестр", "2-ой семестр", "3-ий семестр", "4-ый семестр",
        "5-ый семестр", "6-ой семестр", "7-ой семестр", "8-ой семестр", "ИТОГО"
    ]

    # Инициализация таблицы с заголовками
    table_data = [[(header, 1, 1) for header in headers]]  # Заголовки столбцов

    total_audit_hours = 0
    total_other_practice = 0

    # Заполнение данных по семестрам и подсчет суммы
    for i, row in enumerate(data):
        row_data = [(rows[i], 1, 1)]  # Первый столбец - название строки

        audit_hours = row[0] if len(row) > 0 and row[0] != '' else '-'
        other_practice = row[1] if len(row) > 1 and row[1] != '' else '-'

        if audit_hours != '-' and audit_hours != '':
            total_audit_hours += int(audit_hours.split()[0])

        other_hours = 0
        if other_practice != '-' and other_practice != '':
            other_practice_items = other_practice.split(',')
            for item in other_practice_items:
                other_hours += int(item.split()[0])
            total_other_practice += other_hours

        row_data.append((str(audit_hours), 1, 1))
        row_data.append((str(other_practice), 1, 1))

        table_data.append(row_data)

    total_row = [("ИТОГО", 1, 1), (str(total_audit_hours), 1, 1), (str(total_other_practice), 1, 1)]
    table_data.append(total_row)

    return table_data

# Пример использования
slide_generator = SlideGenerator()

# Определение данных для слайдов
slides_data = [
    [
        {'type': 'text', 'content': "МИНИСТЕРСТВО НАУКИ И ВЫСШЕГО ОБРАЗОВАНИЯ РОССИЙСКОЙ ФЕДЕРАЦИИ\nФЕДЕРАЛЬНОЕ ГОСУДАРСТВЕННОЕ АВТОНОМНОЕ ОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ высшего образования «Национальный исследовательский ядерный университет «МИФИ» (НИЯУ МИФИ)", 'left': Pt(50), 'top': Pt(50), 'width': Pt(620), 'height': Pt(60), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER},
        {'type': 'text', 'content': "ИНСТИТУТ ИНТЕЛЛЕКТУАЛЬНЫХ КИБЕРНЕТИЧЕСКИХ СИСТЕМ\nКАФЕДРА 22 (Кибернетика)", 'left': Pt(50), 'top': Pt(130), 'width': Pt(620), 'height': Pt(40), 'font_size': 12, 'bold': True, 'alignment': PP_ALIGN.CENTER},
        {'type': 'text', 'content': f"НАПРАВЛЕНИЕ ПОДГОТОВКИ\n{education_direction}", 'left': Pt(50), 'top': Pt(190), 'width': Pt(620), 'height': Pt(40), 'font_size': 12, 'bold': True, 'alignment': PP_ALIGN.CENTER},
        {'type': 'text', 'content': f"ПРОФИЛЬ ПОДГОТОВКИ\n{education_profile}", 'left': Pt(50), 'top': Pt(250), 'width': Pt(620), 'height': Pt(40), 'font_size': 12, 'bold': True, 'alignment': PP_ALIGN.CENTER},
        {'type': 'text', 'content': f"Отчет аспиранта за {current_semester} семестр", 'left': Pt(50), 'top': Pt(310), 'width': Pt(620), 'height': Pt(60), 'font_size': 30, 'bold': False, 'alignment': PP_ALIGN.CENTER},
        {'type': 'text', 'content': f"Аспирант: {full_name}\nНаучный руководитель: {supervisor_name}", 'left': Pt(50), 'top': Pt(390), 'width': Pt(620), 'height': Pt(40), 'font_size': 20, 'bold': False, 'alignment': PP_ALIGN.RIGHT},
        {'type': 'text', 'content': "Москва, 2024", 'left': Pt(50), 'top': Pt(450), 'width': Pt(620), 'height': Pt(30), 'font_size': 16, 'bold': False, 'alignment': PP_ALIGN.CENTER},
        {'type': 'text', 'content': "1", 'left': Pt(50), 'top': Pt(490), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (0, 73, 155)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    [
        # {'type': 'text', 'content': "Цели и задачи практики", 'left': Pt(50), 'top': Pt(50), 'width': Pt(620), 'height': Pt(40), 'font_size': 24, 'bold': True, 'alignment': PP_ALIGN.LEFT, 'color': (0, 112, 192)},
        {'type': 'text', 'content': "Объект исследования: " + scientific_obj, 'left': Pt(50), 'top': Pt(100), 'width': Pt(620), 'height': Pt(40), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        {'type': 'text', 'content': "Предмет исследования: " + scientific_subj, 'left': Pt(50), 'top': Pt(150), 'width': Pt(620), 'height': Pt(40), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        {'type': 'table', 'data': [
            [("Общие сведения на текущий момент", 1, 4)],  # Объединение всех 4 столбцов
            [("Руководитель:", 1, 1), (supervisor_name, 1, 1), ("Семестр обучения:", 1, 1), (training_year_fgos, 1, 1)],
            [("Категория:", 1, 1), (category, 1, 1), ("Специальность:", 1, 1), (specialty, 1, 1)],
            [("Кандидатские экзамены:", 3, 1), (candidate_exams, 3, 1), ("Дата поступления:", 1, 1), (enrollment_date, 1, 1)],
            [("", 1, 1), ("", 1, 1), ("", 1, 1), ("", 1, 1)],  # Пустые строки под "Кандидатские экзамены:"
            [("", 1, 1), ("", 1, 1), ("", 1, 1), ("", 1, 1)],  # Пустые строки под "Кандидатские экзамены:"
            [("Тема:", 1, 1), (topic, 1, 3)]  # Объединение трех столбцов для темы
        ], 'left': Inches(1), 'top': Inches(3), 'width': Inches(8), 'height': Inches(2), 'font_size': 12},
        # {'type': 'text', 'content': "Объект исследования: " + scientific_obj, 'left': Pt(350), 'top': Pt(100), 'width': Pt(620), 'height': Pt(100), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        # {'type': 'text', 'content': "Предмет исследования: " + scientific_subj, 'left': Pt(400), 'top': Pt(100), 'width': Pt(620), 'height': Pt(100), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        {'type': 'text', 'content': "Оценка научного руководителя: " + mentor_rate, 'left': Pt(50), 'top': Pt(450), 'width': Pt(620), 'height': Pt(100), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        {'type': 'text', 'content': "2", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    [
        {'type': 'text', 'content': "Отчет о работе над диссертацией за отчетный период", 'left': Pt(50), 'top': Pt(50), 'width': Pt(620), 'height': Pt(40), 'font_size': 24, 'bold': True, 'alignment': PP_ALIGN.LEFT, 'color': (0, 112, 192)},
        {'type': 'text', 'content': f"{report_period_work}\nПроцент выполнения: {progress_percents[current_semester-1]}%\nТема: {topic}\nОбъект исследования: {scientific_obj}\nПредмет исследования: {scientific_subj}\nОценка научного руководителя: {mentor_rate}", 'left': Pt(50), 'top': Pt(100), 'width': Pt(620), 'height': Pt(100), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        {'type': 'text', 'content': "3", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    [
        {'type': 'text', 'content': "Работа над диссертацией (весь период)", 'left': Pt(50), 'top': Pt(50), 'width': Pt(620), 'height': Pt(40), 'font_size': 24, 'bold': True, 'alignment': PP_ALIGN.LEFT, 'color': (0, 112, 192)},
        {'type': 'table', 'data': build_table_disser_work_data(topic, progress_percents, progress_descriptions), 'left': Inches(1), 'top': Inches(1.5), 'width': Inches(8), 'height': Inches(1.5), 'font_size': 12},
        {'type': 'text', 'content': "4", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    [
        {'type': 'text', 'content': "Публикации в журналах за отчетный период", 'left': Pt(50), 'top': Pt(50), 'width': Pt(620), 'height': Pt(40), 'font_size': 24, 'bold': True, 'alignment': PP_ALIGN.LEFT, 'color': (0, 112, 192)},
        {'type': 'table', 'data': build_table_publications_data(publications)[0], 'left': Inches(0.5), 'top': Inches(1.1), 'width': Inches(7), 'height': Inches(1.5), 'font_size': 12, 'column_widths': build_table_publications_data(publications)[1]},
        {'type': 'text', 'content': "5", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    [
        {'type': 'text', 'content': "Участие в конференциях за отчетный период", 'left': Pt(25), 'top': Pt(25), 'width': Pt(620), 'height': Pt(40), 'font_size': 20, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'color': (0, 112, 192)},
        {'type': 'text', 'content': "пока недоступно :(", 'left': Pt(25), 'top': Pt(25), 'width': Pt(620), 'height': Pt(40), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (0, 0, 0)},
        # {'type': 'table', 'data': build_table_publ_all(all_publications), 'left': Inches(0.5), 'top': Inches(1.1), 'width': Inches(8.5), 'height': Inches(1.5), 'font_size': 12},
        {'type': 'text', 'content': "6", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    [
        {'type': 'text', 'content': "Публикации в журналах/участие в конференциях (весь период)", 'left': Pt(25), 'top': Pt(25), 'width': Pt(620), 'height': Pt(40), 'font_size': 20, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'color': (0, 112, 192)},
        {'type': 'table', 'data': build_table_publ_all(all_publications), 'left': Inches(0.5), 'top': Inches(1.1), 'width': Inches(8.5), 'height': Inches(1.5), 'font_size': 12},
        {'type': 'text', 'content': "7", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    [
        {'type': 'text', 'content': "Педагогическая нагрузка за отчетный период", 'left': Pt(25), 'top': Pt(25), 'width': Pt(620), 'height': Pt(40), 'font_size': 24, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'color': (0, 112, 192)},
        {'type': 'table', 'data': build_table_pedagogical_work(pedagogical_data), 'left': Inches(1), 'top': Inches(1.1), 'width': Inches(8), 'height': Inches(4), 'font_size': 12},
        {'type': 'text', 'content': "Другие достижения за отчетный период", 'left': Pt(25), 'top': Pt(385), 'width': Pt(620), 'height': Pt(40), 'font_size': 22, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'color': (0, 112, 192)},
        {'type': 'text', 'content': report_other_achievments, 'left': Pt(50), 'top': Pt(415), 'width': Pt(620), 'height': Pt(100), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        {'type': 'text', 'content': "8", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ],
    # [
    #     {'type': 'text', 'content': "Педагогическая нагрузка (весь период)", 'left': Pt(25), 'top': Pt(25), 'width': Pt(620), 'height': Pt(40), 'font_size': 24, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'color': (0, 112, 192)},
    #     {'type': 'table', 'data': build_table_pedagogical_work_all(pedagogical_data_all), 'left': Inches(0.5), 'top': Inches(1.0), 'width': Inches(8.5), 'height': Inches(1.5), 'font_size': 12},
    #     {'type': 'text', 'content': "9", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
    #     {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
    #     {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    # ],
    [
        {'type': 'text', 'content': "План работы на следующий семестр", 'left': Pt(50), 'top': Pt(50), 'width': Pt(620), 'height': Pt(40), 'font_size': 24, 'bold': True, 'alignment': PP_ALIGN.LEFT, 'color': (0, 112, 192)},
        {'type': 'text', 'content': "\n".join(next_semester_plan), 'left': Pt(50), 'top': Pt(100), 'width': Pt(620), 'height': Pt(200), 'font_size': 14, 'bold': False, 'alignment': PP_ALIGN.LEFT},
        {'type': 'text', 'content': "9", 'left': Pt(50), 'top': Pt(520), 'width': Pt(620), 'height': Pt(20), 'font_size': 12, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'color': (255, 255, 255)},
        {'type': 'image', 'image_path': 'MEPhI_Logo2014_en.png', 'left': Inches(0), 'top': Inches(6.7), 'width': Inches(0.8), 'height': Inches(0.8)},
        {'type': 'image', 'image_path': 'kaf22.png', 'left': Inches(8.5), 'top': Inches(6.9), 'width': Inches(1.5), 'height': Inches(0.6)}
    ]
]

# Создание презентации
slide_generator.create_presentation(slides_data)

# Сохранение презентации
slide_generator.save_presentation('report.pptx')